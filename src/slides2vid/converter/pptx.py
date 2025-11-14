from slides2vid.converter.audio import AudioConverter
from . import  ConverterResult
from pptx import Presentation
from pathlib import Path
from slides2vid.tts.base import TTSEngine


class PPTXAudioConverter(AudioConverter):

    def __init__(self,work_path:Path,pptx_path:Path,engine:TTSEngine,verbose=False) -> None:
        super().__init__(work_path,engine,verbose=verbose)
        self.pptx_path = pptx_path

    def get_slides_text(self)->dict[int,str]:
        presentation = Presentation(str(self.pptx_path))
        texts = {}
        for i,slide in enumerate(presentation.slides):
            if slide.has_notes_slide:
                text = slide.notes_slide.notes_text_frame.text
            else:
                text = ""
            texts[i]=text
        return texts        

    def run(self)-> ConverterResult:
        texts = self.get_slides_text()
        return self.generate_audios(texts)
    