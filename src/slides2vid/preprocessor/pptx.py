from slides2vid.preprocessor.audio import AudioPreprocessor
from slides2vid.preprocessor.core import PreprocessorResult
from pptx import Presentation
from pathlib import Path
from slides2vid.tts.base import TTSEngine


class PPTXAudioPreprocessor(AudioPreprocessor):

    def __init__(self,work_path:Path,engine:TTSEngine,pptx_path:Path) -> None:
        super().__init__(work_path,engine)
        self.pptx_path = pptx_path

    def get_slides_text(self)->dict[int,str]:
        presentation = Presentation(str(self.pptx_path))
        texts = {}
        for i,slide in enumerate(presentation.slides):
            if slide.has_notes_slide:
                text = slide.notes_slide.notes_text_frame.text
            else:
                text = ""
            texts[i]=text
        return texts        
        
    def run(self)-> PreprocessorResult:
        texts = self.get_slides_text()
        changed = self.get_changed(texts)
        result = self.generate_audios(texts,changed)
        self.update_texts_cache(texts)
        return result