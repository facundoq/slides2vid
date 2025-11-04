import faulthandler
from fileinput import filename
from tkinter import LAST

import tqdm
from pathlib import Path

from subprocess import run

from pdf2image import convert_from_path,pdfinfo_from_path
from pptx import Presentation
from gtts import gTTS
from slides2vid.slides import FFMPEGSlideGenerator, SlideGenerator

# from yamlable import YamlAble, yaml_info
import yaml
import logging
logger = logging.getLogger(__name__)

FFMPEG_NAME = 'ffmpeg'

class Project:
    project_filename = "project.yaml"
    TEXT_KEY = "texts"
    LAST_MODIFIED_PDF_KEY = "last_modified_pdf"

    @property
    def project_config_path(self):
        return self.work_path / self.project_filename
    
    def image_path(self,i:int):
        return self.work_path / f'frame_{i}.jpg'
    def audio_path(self,i:int):
        return self.work_path / f'frame_{i}.mp3'
    def video_path(self,i:int):
        return self.work_path / f'frame_{i}.mp4'
    
    def __init__(self, work_folder:Path,language="en"):
        self.work_path = work_folder
        self.texts:dict[int,str] = {}
        self.last_modified_pdf = 0
        self.language = language

        if self.project_config_path.exists():
            with open(self.project_config_path,'r') as f:
                config = yaml.safe_load(f)  
                self.texts = config.get(self.TEXT_KEY, {})
                self.last_modified_pdf = config.get(self.LAST_MODIFIED_PDF_KEY, {})

    def set_text(self, slide_number:int, text:str):
        self.texts[slide_number] = text

    def set_last_modified_pdf(self,date):
        self.last_modified_pdf = date

    def save_project(self):
        config = {self.TEXT_KEY: self.texts, 
                  self.LAST_MODIFIED_PDF_KEY: self.last_modified_pdf}
        with open(self.project_config_path, 'w') as file:
            yaml.dump(config, file)
        

    def generate_audios(self,presentation:Presentation,n:int):
        audio_paths = [self.audio_path(i) for i in range(n)]
        audios_changed = [True for i in range(n)]
        pbar = tqdm.tqdm(enumerate(presentation.slides),total=n)
        pbar.set_description("Generating slide audios")
        for i,slide in pbar:
            if slide.has_notes_slide:
                text = slide.notes_slide.notes_text_frame.text
            else:
                text = "¤"*10
            logger.info(f"Slide {i} text:\n {text}\n")
            audio_changed = True
            if i in self.texts:
                previous_notes = self.texts[i]
                audio_path = audio_paths[i]
                if previous_notes == text and audio_path.exists():
                    logger.info(f"Slide {i} text did not change, skipping audio generation.")
                    audio_changed = False
                    audios_changed[i] = audio_changed
            if audio_changed:
                tts = gTTS(text,lang=self.language)
                tts.save(audio_paths[i])
            self.set_text(i,text)
            return audio_paths,audios_changed
            
                
    def generate_videos(self,generator:SlideGenerator,audio_paths:list[Path],image_paths:list[Path],audio_changed:list[bool],image_changed:list[bool]):
        n = len(audio_paths)
        pbar = tqdm.trange(n)
        video_paths = [self.video_path(i) for i in range(n)]
        pbar.set_description("Generating slide videos")
        for i in pbar:
            if audio_changed[i] or image_changed[i]:
                generator.generate(image_paths[i], audio_paths[i],video_paths[i])
        return video_paths

    def generate_images(self,pdf_path:Path,n:int):
        
        last_modified_pdf = pdf_path.stat().st_mtime
        image_paths = [self.image_path(i) for i in range(n)]
        images_exist = all(map(lambda p: p.exists(),image_paths))
        if last_modified_pdf != self.last_modified_pdf or not images_exist:
            images = convert_from_path(pdf_path)
            pbar = tqdm.tqdm(enumerate(images),total=n)
            pbar.set_description("Generating slide images")
            for i,image in pbar:
                image.save(self.image_path(i))
            images_changed = True
        else:
            images_changed = False

        self.last_modified_pdf = last_modified_pdf
        return image_paths,images_changed
        
    def make_video(self, pptx_path:Path, pdf_path:Path, output_path:Path):
        generator = FFMPEGSlideGenerator(self.work_path)
        pdf_info = pdfinfo_from_path(str(pdf_path))
        print(pdf_info)
        presentation = Presentation(str(pptx_path))
        n = len(presentation.slides)
        assert pdf_info["Pages"] == n
        image_paths, images_changed = self.generate_images(pdf_path,n)
        audio_paths, audio_changed = self.generate_audios(presentation,n)
        video_paths = self.generate_videos(generator,audio_paths,image_paths,audio_changed,images_changed)
        self.save_project()
        logger.info("Concatenating slide videos...")
        generator.concatenate(video_paths,output_path)
        logger.info(" done.")

