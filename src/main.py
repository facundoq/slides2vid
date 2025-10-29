#!/usr/bin/env python

import os
from pathlib import Path
import sys
import tempfile
import argparse
from subprocess import run

from pdf2image import convert_from_path
from pptx import Presentation
from gtts import gTTS
from video import FFMPEGSlideGenerator, SlideGenerator




## Sometimes ffmpeg is avconv
FFMPEG_NAME = 'ffmpeg'
#FFMPEG_NAME = 'avconv'


def ppt_presenter(pptx_path:Path, pdf_path:Path, output_path:Path,work_path:Path):
    generator = FFMPEGSlideGenerator(work_path)
    images_from_path = convert_from_path(pdf_path)
    prs = Presentation(pptx_path)
    assert len(images_from_path) == len(prs.slides)
    video_list = []
    for i, (slide, image) in enumerate(zip(prs.slides, images_from_path)):
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            image_path = work_path / f'frame_{i}.jpg'
            image.save(image_path)
            
            tts = gTTS(text=notes, lang='en')
            audio_path = work_path / f'frame_{i}.mp3'
            tts.save(audio_path)
            video_path = work_path / f'frame_{i}.mp4'
            generator.generate(image_path, audio_path,video_path)
            video_list.append(video_path)
    video_list_str = 'concat:' + '|'.join(video_list)
    ffmpeg_concat(video_list_str, output_path)


def ffmpeg_call(image_path:Path, audio_path:Path, video_path:Path, generator:SlideGenerator,i:int):
    
    
    generator.generate(image_path, audio_path,video_path)
    
    # gen_command = [FFMPEG_NAME, '-loop', '1', '-y', '-i', image_path, '-i', audio_path,
    #       '-c:v', 'libx264', '-tune', 'stillimage', '-c:a', 'aac',
    #       '-b:a', '192k', '-pix_fmt', 'yuv420p', '-shortest', out_path_mp4]
    # print("Running: "+ (" ".join(gen_command)))
    # out = run(gen_command)
    # assert out.returncode == 0, f"Slide {i} generation failed"
    # out = run([FFMPEG_NAME, '-y', '-i', out_path_mp4, '-c', 'copy',
    #       '-bsf:v', 'h264_mp4toannexb', '-f', 'mpegts', out_path_ts])
    # assert out.returncode == 0, f"Slide {i} conversion failed"


def ffmpeg_concat(video_list_str, out_path):
    out =  run([FFMPEG_NAME, '-y', '-f', 'mpegts', '-i', '{}'.format(video_list_str), '-c', 'copy', '-bsf:a', 'aac_adtstoasc', out_path])
    assert out.returncode == 0, f"Slide video concatenation failed"


def main():
    parser = argparse.ArgumentParser(description='PPT Presenter help.')
    parser.add_argument('--pptx', help='input pptx path')
    parser.add_argument('--pdf', help='input pdf path')
    parser.add_argument('-o', '--output', help='output path')
    args = parser.parse_args()
    with tempfile.TemporaryDirectory() as work_path:
        work_path = Path("output")
        os.makedirs(work_path, exist_ok=True)
        ppt_presenter(Path(args.pptx), Path(args.pdf), Path(args.output),work_path)


if __name__ == '__main__':
    main()